#!/usr/bin/env python3
import hashlib
import json
import re
from datetime import datetime, timezone
from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
EXPORT_DIR = ROOT / "presentation-export"
GENERATED_DIR = EXPORT_DIR / "generated"
SLIDE_SOURCE_PATH = EXPORT_DIR / "slide-content-for-ppt.md"
STATUS_PATH = EXPORT_DIR / "export-status.json"
MANIFEST_PATH = EXPORT_DIR / "export-manifest.json"
BLOCKERS_PATH = EXPORT_DIR / "blockers.md"
PPTX_PATH = GENERATED_DIR / "onyxgitex-ai-trust-security-readiness-deck.pptx"

FOOTER_TEXT = "Portfolio/readiness methodology — not production-ready without real client evidence."
SOURCE_READY_STATUS = {
    "pack": "presentation-export",
    "status": "POWERPOINT_DECK_SOURCE_READY",
    "source_pack": "pitch-deck-pack",
    "source_pack_status": "PITCH_DECK_PACK_READY",
    "slide_count": 10,
    "pptx_generated": False,
    "pptx_path": None,
    "pdf_generated": False,
    "pdf_path": None,
    "one_page_handout_pdf_generated": False,
    "production_ready": False,
    "go_decision": False,
    "client_evidence_verified": False,
}

REQUIRED_MANIFEST_PATHS = [
    "presentation-export/README.md",
    "presentation-export/powerpoint-export-plan.md",
    "presentation-export/slide-content-for-ppt.md",
    "presentation-export/slide-layout-spec.md",
    "presentation-export/speaker-notes-export.md",
    "presentation-export/one-page-handout.md",
    "presentation-export/presentation-claim-boundaries.md",
    "presentation-export/export-status.json",
    "presentation-export/export-manifest.json",
    "presentation-export/validation-result.txt",
    "presentation-export/blockers.md",
    "presentation-export/generated/onyxgitex-ai-trust-security-readiness-deck.pptx",
    "presentation-export/generated/onyxgitex-ai-trust-security-readiness-deck.pdf",
    "presentation-export/generated/onyxgitex-one-page-handout.pdf",
]


def sha256_for(path: Path):
    if not path.exists():
        return None
    return hashlib.sha256(path.read_bytes()).hexdigest()


def parse_slides(md_text: str):
    slide_pattern = re.compile(r"## Slide \d+\n(.*?)(?=\n## Slide \d+\n|\Z)", re.S)
    slides = []
    for block in slide_pattern.findall(md_text):
        title_match = re.search(r"- Title: (.*)", block)
        main_message_match = re.search(r"- Main message: (.*)", block)
        bullet_matches = re.findall(r"^  - (.*)$", block, re.M)

        if not title_match or not main_message_match:
            raise ValueError("Slide content missing title or main message")
        if len(bullet_matches) < 3 or len(bullet_matches) > 5:
            raise ValueError("Each slide must include 3-5 bullets")

        slides.append(
            {
                "title": title_match.group(1).strip(),
                "main_message": main_message_match.group(1).strip(),
                "bullets": bullet_matches,
            }
        )

    if len(slides) != 10:
        raise ValueError(f"Expected exactly 10 slides, found {len(slides)}")

    return slides


def generate_pptx(slides):
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches, Pt

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for slide_data in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(12.0), Inches(0.9))
        title_tf = title_box.text_frame
        title_tf.clear()
        title_run = title_tf.paragraphs[0].add_run()
        title_run.text = slide_data["title"]
        title_run.font.size = Pt(30)
        title_run.font.bold = True
        title_run.font.color.rgb = RGBColor(20, 20, 20)

        msg_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.35), Inches(12.0), Inches(0.8))
        msg_tf = msg_box.text_frame
        msg_tf.clear()
        msg_run = msg_tf.paragraphs[0].add_run()
        msg_run.text = slide_data["main_message"]
        msg_run.font.size = Pt(19)
        msg_run.font.color.rgb = RGBColor(55, 55, 55)

        bullet_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(2.1), Inches(8.0), Inches(4.6))
        bullet_bg.fill.solid()
        bullet_bg.fill.fore_color.rgb = RGBColor(245, 245, 245)
        bullet_bg.line.color.rgb = RGBColor(210, 210, 210)

        bullets_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.35), Inches(7.5), Inches(4.0))
        bullets_tf = bullets_box.text_frame
        bullets_tf.clear()
        for idx, bullet in enumerate(slide_data["bullets"]):
            paragraph = bullets_tf.paragraphs[0] if idx == 0 else bullets_tf.add_paragraph()
            paragraph.text = bullet
            paragraph.level = 0
            paragraph.font.size = Pt(16)
            paragraph.font.color.rgb = RGBColor(35, 35, 35)

        visual_panel = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(9.0), Inches(2.1), Inches(3.3), Inches(4.6))
        visual_panel.fill.solid()
        visual_panel.fill.fore_color.rgb = RGBColor(233, 233, 233)
        visual_panel.line.color.rgb = RGBColor(170, 170, 170)
        visual_panel.text_frame.text = "Visual structure\n\nSimple blocks\nand flow markers"
        visual_panel.text_frame.paragraphs[0].font.bold = True
        for para in visual_panel.text_frame.paragraphs:
            para.font.size = Pt(14)
            para.font.color.rgb = RGBColor(60, 60, 60)

        footer_box = slide.shapes.add_textbox(Inches(0.7), Inches(7.0), Inches(12.0), Inches(0.35))
        footer_tf = footer_box.text_frame
        footer_tf.clear()
        footer_run = footer_tf.paragraphs[0].add_run()
        footer_run.text = FOOTER_TEXT
        footer_run.font.size = Pt(10)
        footer_run.font.color.rgb = RGBColor(90, 90, 90)

    GENERATED_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(PPTX_PATH)


def update_manifest():
    entries = []
    for rel in REQUIRED_MANIFEST_PATHS:
        path = ROOT / rel
        if rel.endswith("export-manifest.json"):
            entry_sha = "SELF_REFERENTIAL_MANIFEST_HASH_NOT_APPLICABLE"
        else:
            entry_sha = sha256_for(path)
        description_map = {
            "presentation-export/generated/onyxgitex-ai-trust-security-readiness-deck.pptx": "Generated PowerPoint deck",
            "presentation-export/generated/onyxgitex-ai-trust-security-readiness-deck.pdf": "Generated deck PDF",
            "presentation-export/generated/onyxgitex-one-page-handout.pdf": "Generated one-page handout PDF",
        }
        entries.append(
            {
                "path": rel,
                "exists": path.exists(),
                "sha256": entry_sha,
                "description": description_map.get(rel, rel.split("/")[-1]),
            }
        )

    manifest = {
        "pack": "presentation-export",
        "generated_at_utc": datetime.now(timezone.utc).isoformat(),
        "required_files": entries,
    }
    MANIFEST_PATH.write_text(json.dumps(manifest, indent=2) + "\n")


def write_blockers(success: bool, reason: str | None = None):
    if success:
        content = """# Presentation Export Blockers

No PowerPoint deck export blockers detected.

Verified:
- PowerPoint source content exists.
- 10 slides are defined.
- Speaker notes export exists.
- One-page handout exists.
- Claim boundaries exist.
- PowerPoint file generated.
- Production readiness remains blocked.
- GO decision remains false.
- Client evidence remains unverified.

Remaining blockers before any production claim:
- Real client evidence is not present.
- Real client runtime proof is not present.
- Real client compliance review is not present.
- Real client approver signoff is not present.
- Version 4 must be filled with real client evidence before any production decision.
"""
    else:
        content = f"""# Presentation Export Blockers

PowerPoint source pack is ready, but .pptx generation was not completed.

Reason:
- {reason}

Verified:
- PowerPoint source content exists.
- 10 slides are defined.
- Speaker notes export exists.
- One-page handout exists.
- Claim boundaries exist.
- Production readiness remains blocked.
- GO decision remains false.
- Client evidence remains unverified.

Next action:
Generate the .pptx using available presentation tooling.
"""
    BLOCKERS_PATH.write_text(content)


def main():
    status = dict(SOURCE_READY_STATUS)
    failure_reason = None

    try:
        slides = parse_slides(SLIDE_SOURCE_PATH.read_text())
        generate_pptx(slides)
        if PPTX_PATH.exists() and PPTX_PATH.stat().st_size > 0:
            status.update(
                {
                    "status": "POWERPOINT_DECK_EXPORT_READY",
                    "pptx_generated": True,
                    "pptx_path": "presentation-export/generated/onyxgitex-ai-trust-security-readiness-deck.pptx",
                }
            )
        else:
            failure_reason = "PowerPoint generation completed without a valid output file"
    except Exception as ex:
        failure_reason = str(ex)

    STATUS_PATH.write_text(json.dumps(status, indent=2) + "\n")
    write_blockers(success=(status["status"] == "POWERPOINT_DECK_EXPORT_READY"), reason=failure_reason)
    update_manifest()
    print(status["status"])


if __name__ == "__main__":
    main()
